from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUT = Path(__file__).with_name("在线编程助手算法设计_两页摘要.pptx")

SLIDE_W = 13.333
SLIDE_H = 7.5

BG = "07111F"
PANEL = "0D1B2E"
PANEL_2 = "10243A"
WHITE = "F6F9FC"
MUTED = "9DB0C6"
CYAN = "35D3FF"
CYAN_DARK = "123B50"
GREEN = "50E3A4"
GREEN_DARK = "123D35"
AMBER = "FFC857"
AMBER_DARK = "4A3919"
RED = "FF6B7A"
RED_DARK = "4B202B"
LINE = "27405B"

FONT = "Noto Sans CJK SC"
MONO = "DejaVu Sans Mono"


def rgb(hex_color: str) -> RGBColor:
    return RGBColor.from_string(hex_color)


def set_fill(shape, color: str, transparency: int = 0):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(color)
    shape.fill.transparency = transparency


def set_line(shape, color: str, width: float = 1.0, transparency: int = 0):
    shape.line.color.rgb = rgb(color)
    shape.line.width = Pt(width)
    shape.line.transparency = transparency


def add_text(
    slide,
    text,
    x,
    y,
    w,
    h,
    size=14,
    color=WHITE,
    bold=False,
    font=FONT,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.MIDDLE,
    margin=0,
    line_spacing=1.0,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    box.text_frame.clear()
    box.text_frame.margin_left = Inches(margin)
    box.text_frame.margin_right = Inches(margin)
    box.text_frame.margin_top = Inches(margin)
    box.text_frame.margin_bottom = Inches(margin)
    box.text_frame.vertical_anchor = valign
    p = box.text_frame.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    return box


def add_rich_text(
    slide,
    parts,
    x,
    y,
    w,
    h,
    size=14,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.MIDDLE,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    box.text_frame.clear()
    box.text_frame.margin_left = 0
    box.text_frame.margin_right = 0
    box.text_frame.margin_top = 0
    box.text_frame.margin_bottom = 0
    box.text_frame.vertical_anchor = valign
    p = box.text_frame.paragraphs[0]
    p.alignment = align
    for part in parts:
        run = p.add_run()
        run.text = part["text"]
        run.font.name = part.get("font", FONT)
        run.font.size = Pt(part.get("size", size))
        run.font.bold = part.get("bold", False)
        run.font.color.rgb = rgb(part.get("color", WHITE))
    return box


def add_round_rect(slide, x, y, w, h, fill=PANEL, line=LINE, radius=True):
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    set_fill(shape, fill)
    set_line(shape, line, 1.0)
    if radius:
        try:
            shape.adjustments[0] = 0.08
        except (IndexError, ValueError):
            pass
    return shape


def add_pill(slide, text, x, y, w, h, fill, color=WHITE, size=10, line=None):
    shape = add_round_rect(slide, x, y, w, h, fill=fill, line=line or fill)
    add_text(
        slide,
        text,
        x,
        y,
        w,
        h,
        size=size,
        color=color,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    return shape


def add_circle_label(slide, text, x, y, d, fill, color=BG, size=11):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d)
    )
    set_fill(shape, fill)
    set_line(shape, fill, 0.7)
    add_text(
        slide,
        text,
        x,
        y,
        d,
        d,
        size=size,
        color=color,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    return shape


def add_chevron(slide, x, y, w=0.22, h=0.34, color=LINE):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.CHEVRON, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    set_fill(shape, color)
    shape.line.fill.background()
    return shape


def add_hline(slide, x1, y1, x2, y2, color=LINE, width=1.2):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1),
        Inches(y1),
        Inches(x2),
        Inches(y2),
    )
    set_line(line, color, width)
    return line


def add_header(slide, eyebrow, title, page):
    add_text(
        slide,
        eyebrow.upper(),
        0.62,
        0.30,
        5.6,
        0.25,
        size=9,
        color=CYAN,
        bold=True,
    )
    add_text(slide, title, 0.62, 0.58, 11.6, 0.54, size=25, bold=True)
    add_text(
        slide,
        f"0{page}",
        12.12,
        0.38,
        0.55,
        0.34,
        size=11,
        color=MUTED,
        bold=True,
        align=PP_ALIGN.RIGHT,
    )
    add_hline(slide, 12.78, 0.55, 13.0, 0.55, color=CYAN, width=2.4)


def add_footer(slide, source="来源：《在线编程助手算法设计》"):
    add_text(slide, source, 0.62, 7.13, 5.6, 0.18, size=8, color=MUTED)
    add_text(
        slide,
        "TREE-SITTER × VIRTUAL INTERPRETER",
        8.0,
        7.13,
        4.7,
        0.18,
        size=8,
        color=MUTED,
        bold=True,
        align=PP_ALIGN.RIGHT,
    )


def add_stage_card(
    slide,
    x,
    y,
    w,
    h,
    number,
    kicker,
    title,
    body,
    accent,
    accent_dark,
):
    add_round_rect(slide, x, y, w, h, fill=PANEL, line=LINE)
    # Accent rail
    rail = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.055), Inches(h)
    )
    set_fill(rail, accent)
    rail.line.fill.background()
    add_circle_label(slide, number, x + 0.18, y + 0.19, 0.34, accent, color=BG, size=10)
    add_text(
        slide,
        kicker.upper(),
        x + 0.62,
        y + 0.18,
        w - 0.8,
        0.24,
        size=8.5,
        color=accent,
        bold=True,
    )
    add_text(
        slide,
        title,
        x + 0.18,
        y + 0.68,
        w - 0.36,
        0.54,
        size=17,
        color=WHITE,
        bold=True,
        valign=MSO_ANCHOR.TOP,
    )
    add_text(
        slide,
        body,
        x + 0.18,
        y + 1.35,
        w - 0.36,
        h - 1.56,
        size=10.4,
        color=MUTED,
        valign=MSO_ANCHOR.TOP,
        line_spacing=1.15,
    )
    add_pill(
        slide,
        "CORE",
        x + w - 0.62,
        y + h - 0.36,
        0.42,
        0.2,
        fill=accent_dark,
        color=accent,
        size=7.5,
        line=accent_dark,
    )


def build_slide_1(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    background.fill.solid()
    background.fill.fore_color.rgb = rgb(BG)

    add_header(
        slide,
        "Algorithm overview / 架构总览",
        "在线编程助手：把自然语言变成可信机器人代码",
        1,
    )

    add_rich_text(
        slide,
        [
            {"text": "核心思路：", "color": MUTED, "size": 12},
            {"text": "Tree-sitter", "color": CYAN, "bold": True, "size": 12},
            {"text": " 保证代码“写得对”，", "color": MUTED, "size": 12},
            {"text": "虚拟解释器", "color": GREEN, "bold": True, "size": 12},
            {"text": " 保证代码“跑得通”。", "color": MUTED, "size": 12},
        ],
        0.62,
        1.22,
        9.2,
        0.34,
    )
    add_pill(
        slide,
        "Siasun / KRL",
        10.72,
        1.20,
        1.34,
        0.34,
        fill=CYAN_DARK,
        color=CYAN,
        size=9,
    )
    add_pill(
        slide,
        "AST 驱动",
        12.14,
        1.20,
        0.78,
        0.34,
        fill=GREEN_DARK,
        color=GREEN,
        size=9,
    )

    y = 1.90
    h = 2.78
    xs = [0.50, 3.04, 5.58, 8.12, 10.66]
    widths = [2.22, 2.22, 2.22, 2.22, 2.18]
    cards = [
        (
            "1",
            "INPUT",
            "用户意图",
            "自然语言描述任务\n例如：移动到 P1 并等待",
            AMBER,
            AMBER_DARK,
        ),
        (
            "2",
            "GENERATION",
            "代码草稿",
            "LLM 或模板生成\n遵循 SET / IF / GOTO / MOVJ 等关键词",
            CYAN,
            CYAN_DARK,
        ),
        (
            "3",
            "STATIC GATE",
            "语法与结构",
            "Tree-sitter 解析 AST\n检查 ERROR / MISSING、NOP / END、参数与类型",
            CYAN,
            CYAN_DARK,
        ),
        (
            "4",
            "DYNAMIC GATE",
            "语义与运行",
            "解释器逐行执行 AST\n模拟寄存器、控制流、运动与 IO",
            GREEN,
            GREEN_DARK,
        ),
        (
            "5",
            "DELIVERY",
            "可信输出",
            "代码 + AST 高亮\n最终状态 + 执行日志 / 测试报告",
            GREEN,
            GREEN_DARK,
        ),
    ]
    for i, (x, w, data) in enumerate(zip(xs, widths, cards)):
        add_stage_card(slide, x, y, w, h, *data)
        if i < len(cards) - 1:
            add_chevron(slide, x + w + 0.08, y + 1.20, color=LINE)

    add_pill(
        slide,
        "第一道门禁 · 写得对",
        5.80,
        1.64,
        1.77,
        0.31,
        fill=CYAN_DARK,
        color=CYAN,
        size=9,
    )
    add_pill(
        slide,
        "第二道门禁 · 跑得通",
        8.30,
        1.64,
        1.87,
        0.31,
        fill=GREEN_DARK,
        color=GREEN,
        size=9,
    )

    # Build rail
    add_round_rect(slide, 0.50, 5.07, 12.34, 1.64, fill="091827", line=LINE)
    add_pill(
        slide,
        "一次性初始化",
        0.76,
        5.30,
        1.05,
        0.30,
        fill=PANEL_2,
        color=MUTED,
        size=8.5,
        line=LINE,
    )
    add_text(
        slide,
        "语法库构建链",
        0.76,
        5.72,
        1.45,
        0.33,
        size=13,
        color=WHITE,
        bold=True,
    )
    add_text(
        slide,
        "支撑后续每次解析与高亮",
        0.76,
        6.10,
        1.75,
        0.25,
        size=8.5,
        color=MUTED,
    )

    rail_items = [
        ("grammar.js", "Siasun / KRL 语法定义", CYAN),
        ("Build Library", "编译为 .so / .dll", AMBER),
        ("Parser Load", "Python 加载语言并生成 AST", GREEN),
    ]
    rail_x = [2.76, 6.20, 9.64]
    rail_w = 2.62
    for idx, ((label, desc, accent), x) in enumerate(zip(rail_items, rail_x)):
        add_circle_label(slide, str(idx + 1), x, 5.43, 0.30, accent, size=9)
        add_text(slide, label, x + 0.42, 5.34, 2.0, 0.30, size=12, bold=True)
        add_text(
            slide,
            desc,
            x + 0.42,
            5.76,
            rail_w - 0.42,
            0.33,
            size=9.5,
            color=MUTED,
        )
        if idx < 2:
            add_hline(slide, x + 2.43, 5.78, x + 3.18, 5.78, color=LINE, width=1.5)
            add_chevron(slide, x + 3.06, 5.61, w=0.20, h=0.34, color=LINE)
    add_footer(slide)
    return slide


def add_flow_card(slide, x, y, w, h, step, title, body, accent, accent_dark):
    add_round_rect(slide, x, y, w, h, fill=PANEL, line=LINE)
    add_pill(
        slide,
        f"STEP {step}",
        x + 0.16,
        y + 0.16,
        0.68,
        0.25,
        fill=accent_dark,
        color=accent,
        size=8,
        line=accent_dark,
    )
    add_text(
        slide,
        title,
        x + 0.16,
        y + 0.56,
        w - 0.32,
        0.40,
        size=15,
        bold=True,
    )
    add_text(
        slide,
        body,
        x + 0.16,
        y + 1.02,
        w - 0.32,
        h - 1.16,
        size=9.5,
        color=MUTED,
        valign=MSO_ANCHOR.TOP,
        line_spacing=1.15,
    )


def build_slide_2(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    background.fill.solid()
    background.fill.fore_color.rgb = rgb(BG)

    add_header(
        slide,
        "Validation loop / 验证闭环",
        "从候选代码到可审计结果：失败即定位、反馈、再生成",
        2,
    )
    add_rich_text(
        slide,
        [
            {"text": "静态错误", "color": CYAN, "bold": True, "size": 11},
            {"text": " 与 ", "color": MUTED, "size": 11},
            {"text": "运行时异常", "color": GREEN, "bold": True, "size": 11},
            {"text": " 共用一条反馈通道；超过重试上限则带警告返回最后版本。", "color": MUTED, "size": 11},
        ],
        0.62,
        1.20,
        9.3,
        0.32,
    )
    add_pill(
        slide,
        "MAX_RETRIES = 3",
        10.95,
        1.18,
        1.95,
        0.35,
        fill=RED_DARK,
        color=RED,
        size=9,
    )

    # Main workflow canvas
    add_round_rect(slide, 0.50, 1.72, 8.92, 5.12, fill="091827", line=LINE)
    add_text(
        slide,
        "主处理链",
        0.76,
        1.91,
        1.0,
        0.23,
        size=9,
        color=MUTED,
        bold=True,
    )

    fy = 2.27
    fh = 1.72
    flows = [
        (0.78, 1.72, "1", "生成候选", "用户意图\n→ LLM / 模板", AMBER, AMBER_DARK),
        (
            2.78,
            1.90,
            "2",
            "解析与静态校验",
            "生成 AST\n检查语法与结构",
            CYAN,
            CYAN_DARK,
        ),
        (
            4.98,
            1.90,
            "3",
            "解释器动态执行",
            "更新寄存器 / PC\n模拟运动与 IO",
            GREEN,
            GREEN_DARK,
        ),
        (
            7.18,
            1.90,
            "4",
            "成功交付",
            "代码 + 高亮\n状态 + 日志",
            GREEN,
            GREEN_DARK,
        ),
    ]
    for idx, (x, w, *args) in enumerate(flows):
        add_flow_card(slide, x, fy, w, fh, *args)
        if idx < len(flows) - 1:
            add_chevron(slide, x + w + 0.14, fy + 0.70, w=0.20, h=0.34, color=LINE)

    add_pill(
        slide,
        "ERROR / MISSING · NOP / END · 参数 / 类型",
        2.87,
        4.15,
        1.72,
        0.30,
        fill=CYAN_DARK,
        color=CYAN,
        size=7.8,
    )
    add_pill(
        slide,
        "除零 · 死循环 · 未定义标签 · 断言",
        5.07,
        4.15,
        1.72,
        0.30,
        fill=GREEN_DARK,
        color=GREEN,
        size=7.8,
    )

    # Error branch to feedback loop
    add_hline(slide, 3.73, 3.99, 3.73, 4.64, color=RED, width=1.1)
    add_hline(slide, 5.93, 3.99, 5.93, 4.64, color=RED, width=1.1)
    add_text(
        slide,
        "失败",
        3.86,
        4.52,
        0.42,
        0.22,
        size=8,
        color=RED,
        bold=True,
    )
    add_text(
        slide,
        "异常",
        6.06,
        4.52,
        0.42,
        0.22,
        size=8,
        color=RED,
        bold=True,
    )

    loop_y = 4.78
    loop_h = 1.18
    loop_items = [
        (
            1.17,
            2.02,
            "① 错误定位",
            "行号 · 类型 · 当前代码",
            RED,
            RED_DARK,
        ),
        (
            3.55,
            2.08,
            "② 构建修正 Prompt",
            "目标 + 代码 + error_msg",
            AMBER,
            AMBER_DARK,
        ),
        (
            5.99,
            2.03,
            "③ 判断重试",
            "retry_count < 3 ?",
            CYAN,
            CYAN_DARK,
        ),
    ]
    for idx, (x, w, title, body, accent, accent_dark) in enumerate(loop_items):
        add_round_rect(slide, x, loop_y, w, loop_h, fill=PANEL, line=accent_dark)
        add_text(
            slide,
            title,
            x + 0.15,
            loop_y + 0.15,
            w - 0.3,
            0.32,
            size=11,
            color=accent,
            bold=True,
        )
        add_text(
            slide,
            body,
            x + 0.15,
            loop_y + 0.58,
            w - 0.3,
            0.28,
            size=8.7,
            color=MUTED,
        )
        if idx < len(loop_items) - 1:
            add_chevron(slide, x + w + 0.12, loop_y + 0.42, w=0.18, h=0.32, color=LINE)

    # Feedback return rail
    add_hline(slide, 7.00, 6.16, 1.62, 6.16, color=CYAN, width=1.4)
    add_hline(slide, 1.62, 6.16, 1.62, 6.01, color=CYAN, width=1.4)
    add_chevron(slide, 1.50, 5.87, w=0.24, h=0.34, color=CYAN)
    add_text(
        slide,
        "是 → 返回生成器",
        4.06,
        6.18,
        1.45,
        0.22,
        size=8.5,
        color=CYAN,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_pill(
        slide,
        "否 → 输出最后代码 + 错误警告",
        5.98,
        6.38,
        2.76,
        0.28,
        fill=RED_DARK,
        color=RED,
        size=8.3,
    )

    # Delivery panel
    add_round_rect(slide, 9.70, 1.72, 3.13, 5.12, fill=PANEL, line=LINE)
    add_text(
        slide,
        "通过后交付什么？",
        9.98,
        1.98,
        2.56,
        0.35,
        size=15,
        color=WHITE,
        bold=True,
    )
    add_text(
        slide,
        "四类结果共同构成可审计输出",
        9.98,
        2.36,
        2.56,
        0.28,
        size=9,
        color=MUTED,
    )
    outputs = [
        ("01", "生成代码", "可直接编辑 / 下发"),
        ("02", "AST 高亮", "关键词、变量、结构"),
        ("03", "最终状态", "I / R / P / OT 寄存器"),
        ("04", "执行报告", "轨迹、日志、异常信息"),
    ]
    oy = 2.86
    for num, title, desc in outputs:
        add_circle_label(slide, num, 9.98, oy, 0.34, GREEN, color=BG, size=7.5)
        add_text(slide, title, 10.48, oy - 0.01, 1.55, 0.25, size=11, bold=True)
        add_text(
            slide,
            desc,
            10.48,
            oy + 0.28,
            1.98,
            0.23,
            size=8.2,
            color=MUTED,
        )
        oy += 0.76

    add_pill(
        slide,
        "示例：PASSED",
        9.98,
        5.97,
        1.06,
        0.28,
        fill=GREEN_DARK,
        color=GREEN,
        size=8.2,
    )
    trace = add_round_rect(slide, 9.98, 6.31, 2.56, 0.35, fill="081522", line=LINE)
    add_text(
        slide,
        "NOP → SET I1 5 → OUT ON → END",
        10.08,
        6.31,
        2.36,
        0.35,
        size=7.8,
        color=WHITE,
        font=MONO,
        align=PP_ALIGN.CENTER,
    )

    add_footer(slide)
    return slide


def set_metadata(prs):
    props = prs.core_properties
    props.title = "在线编程助手算法设计｜两页摘要"
    props.subject = "基于 Tree-sitter 与虚拟解释器的机器人编程助手"
    props.author = "OpenAI Codex"
    props.keywords = "Tree-sitter, Siasun, KRL, AST, 虚拟解释器, 在线编程助手"
    props.comments = "根据 scripts/在线编程助手算法设计.md 汇总"


def validate_bounds(prs):
    max_x = prs.slide_width
    max_y = prs.slide_height
    issues = []
    for slide_idx, slide in enumerate(prs.slides, start=1):
        for shape_idx, shape in enumerate(slide.shapes, start=1):
            if shape.left < 0 or shape.top < 0:
                issues.append(f"slide {slide_idx} shape {shape_idx}: negative position")
            if shape.left + shape.width > max_x + Inches(0.01):
                issues.append(f"slide {slide_idx} shape {shape_idx}: exceeds width")
            if shape.top + shape.height > max_y + Inches(0.01):
                issues.append(f"slide {slide_idx} shape {shape_idx}: exceeds height")
    if issues:
        raise ValueError("\n".join(issues))


def main():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    set_metadata(prs)
    build_slide_1(prs)
    build_slide_2(prs)
    validate_bounds(prs)
    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    main()
